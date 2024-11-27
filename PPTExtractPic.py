import pptx
from pathlib import Path
from loguru import logger
import multiprocessing as mp
from pptx.enum.shapes import MSO_SHAPE_TYPE
from tqdm import tqdm


def extract_pic(pptx_file):
    prs = pptx.Presentation(pptx_file)
    for i, slide in tqdm(enumerate(prs.slides), total=len(prs.slides), desc=pptx_file.stem, leave=False):
        for shape in slide.shapes:
            if not shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                continue
            image = shape.image
            image_bytes = image.blob
            image_path = Path('images') / \
                f"{pptx_file.stem}_{i}_{shape.name}.png"
            if not image_path.exists():
                with open(image_path, 'wb') as f:
                    f.write(image_bytes)


if __name__ == '__main__':
    folders = [
        r'C:\Users\xushi\Nextcloud'
    ]
    pptx_files = []
    for folder in folders:
        for file in Path(folder).rglob('*.pptx'):
            pptx_files.append(file)
    logger.info(f"Found {len(pptx_files)} pptx files")

    Path('images').mkdir(exist_ok=True, parents=True)

    for pptx_file in tqdm(pptx_files, desc='Extracting images'):
        try:
            extract_pic(pptx_file)
        except Exception as e:
            logger.error(f"Error processing {pptx_file}: {e}")
