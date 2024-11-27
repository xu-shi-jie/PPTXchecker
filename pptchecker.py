"""Main file for PPTChecker"""

import argparse
import sys
from time import sleep
from pptx import Presentation
from util import display_comments_on_webpage, is_backup_slide, read_config_yaml
from rules import (
    must_end_with_summary_slide,
    should_have_slide_numbers,
    has_smooth_slide_transitions,
    should_have_high_contrast_fonts_colours,
    should_not_have_excessive_text,
    does_not_have_complete_sentences,
    estimate_presentation_length
)
import os
import datetime
from loguru import logger

parser = argparse.ArgumentParser(description='Analyze')
parser.add_argument('-p', '--presentation', type=str,
                    default=r'C:\Users\xushi\Nextcloud\Shijie\PhD Defense\Slides\PhD_Defense_ShijieXu.pptx')
parser.add_argument('-o', '--output', type=str, default="output.html")
args = parser.parse_args()


def main_controller(prs, config):
    slide_feedback = []
    for i, slide in enumerate(prs.slides):
        if i == 0:
            start_slide_num = prs.slides.index(slide)
        if is_backup_slide(slide):
            break
        slide_feedback.append("")

    general_feedback = ""
    pass_all_checks = True

    satisfied = must_end_with_summary_slide(prs)
    if not satisfied:
        general_feedback += "Please end the presention with a summary slide.<br>"

    satisfied = should_have_slide_numbers(prs, slide_feedback)
    if not satisfied:
        general_feedback += "Please add slide numbers.<br>"

    satisfied = has_smooth_slide_transitions(prs, config, slide_feedback)
    if not satisfied:
        general_feedback += "Please check slide transitions.<br>"

    satisfied = should_have_high_contrast_fonts_colours(
        prs, config, slide_feedback)
    if not satisfied:
        general_feedback += "Please check colours and fonts.<br>"

    satisfied = should_not_have_excessive_text(prs, config, slide_feedback)
    if not satisfied:
        general_feedback += "Please ensure that slides do not have too much text.<br>"

    does_not_have_complete_sentences(prs, slide_feedback)

    time_estimate, slide_times, cumul_slide_times = estimate_presentation_length(
        prs, config)
    if time_estimate:
        print("Estimate total time for presentation: ", time_estimate)
    else:
        print("Cannot estimate presentation time without any speaker notes provided!\n")

    for slide_i, _ in enumerate(slide_feedback):
        feedback = slide_feedback[slide_i]
        if feedback:
            pass_all_checks = False
            slide_feedback[slide_i] = feedback.replace('\n', '<br>')

    display_info = {}
    display_info["start_slide_num"] = start_slide_num
    display_info["slide_feedback"] = slide_feedback
    display_info["slide_times"] = slide_times
    display_info["cumul_slide_times"] = cumul_slide_times
    display_info["general_feedback"] = general_feedback

    display_comments_on_webpage(time_estimate, display_info,
                                pass_all_checks, args.output)


def main():
    if not args.presentation:
        print("Must provide a presentation file.")
        sys.exit()
    if not args.presentation.endswith(".pptx"):
        print("Input file must be of '.pptx' type.")
        sys.exit()

    yaml_file = "./config/default.yaml"

    last_modified_yaml, last_modified_pptx = -1, -1

    while True:
        if last_modified_yaml < os.path.getmtime(yaml_file) or last_modified_pptx < os.path.getmtime(args.presentation):
            logger.info(f'pptx file or yaml file changed.')
            # check if pptx file changed or yaml file changed
            config = read_config_yaml(yaml_file)
            path_to_presentation = args.presentation
            prs = Presentation(path_to_presentation)
            main_controller(prs, config)
            last_modified_yaml = os.path.getmtime(yaml_file)
            last_modified_pptx = os.path.getmtime(args.presentation)
        sleep(0.5)


if __name__ == "__main__":
    main()
